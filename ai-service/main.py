from fastapi import FastAPI, UploadFile, File, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from docx import Document
from pptx import Presentation
import io

app = FastAPI(title="Med-Zukkoo AI Microservice", version="1.0")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

def extract_text_from_docx(file_bytes):
    doc = Document(io.BytesIO(file_bytes))
    return "\n".join([para.text for para in doc.paragraphs if para.text.strip()])

def extract_text_from_pdf(file_bytes):
    import fitz # PyMuPDF
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    text = ""
    for page in doc:
        text += page.get_text()
    return text

def extract_text_from_pptx(file_bytes):
    prs = Presentation(io.BytesIO(file_bytes))
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

import os
import json
from groq import Groq
from dotenv import load_dotenv

load_dotenv()
GROQ_API_KEY = os.getenv("GROQ_API_KEY")

# Initialize real Groq client (requires API key in .env file)
# Example .env: GROQ_API_KEY=gsk_your_real_key_here
try:
    groq_client = Groq(api_key=GROQ_API_KEY) if GROQ_API_KEY else None
except Exception:
    groq_client = None

def analyze_with_real_groq(text):
    if not groq_client:
        # Fallback agar Groq ulanmagan yoki API_KEY kiritilmagan bo'lsa
        return {
            "success": False,
            "error": "GROQ_NOT_CONNECTED",
            "message": "Groq API kaliti ulanmagan. .env fayliga GROQ_API_KEY ni kiriting!"
        }

    try:
        completion = groq_client.chat.completions.create(
            model="llama3-8b-8192",
            messages=[
                {
                    "role": "system",
                    "content": "You are a Medical Education AI evaluator. First, determine if the provided text is related to medicine, human anatomy, diseases, or pharmacology. If it is NOT, output exactly a JSON: {\"success\": false, \"message\": \"Not medical context\"}. If it IS medical, output a JSON with generated 'flashcards' and 'quizzes' based on the text. Output ONLY valid JSON."
                },
                {
                    "role": "user",
                    "content": text[:4000] # Pass first 4k chars to respect context limits
                }
            ],
            temperature=0.3,
            response_format={"type": "json_object"}
        )
        response_json = json.loads(completion.choices[0].message.content)
        
        if not response_json.get("success", True):
            return {
                "success": False,
                "error": "TIBBIYOTGA_ALOQADOR_EMAS",
                "message": "Groq AI Llama-3: Kechirasiz, yuklangan fayl tibbiyot yoki farmakologiyaga oid emas. Iltimos, asosan faqat klinik materiallarni yuklang!"
            }
            
        return {
            "success": True,
            "flashcards": response_json.get("flashcards", [{"q": "AI test generated", "a": "Test Data"}]),
            "quizzes": response_json.get("quizzes", [])
        }
    except Exception as e:
        return {
            "success": False,
            "error": "GROQ_API_ERROR",
            "message": f"Groq AI bilan ulanishda xatolik: {str(e)}"
        }

@app.post("/api/ai/process-file")
async def process_educational_file(file: UploadFile = File(...)):
    """
    Accepts DOCX, PDF, or PPTX.
    Extracts text to be sent to the actual Groq LLM for Medical Validation and Quiz/Flashcard generation.
    """
    try:
        content = await file.read()
        filename = file.filename.lower()
        extracted_text = ""

        if filename.endswith(".docx"):
            extracted_text = extract_text_from_docx(content)
        elif filename.endswith(".pdf"):
            extracted_text = extract_text_from_pdf(content)
        elif filename.endswith(".pptx"):
            extracted_text = extract_text_from_pptx(content)
        else:
            raise HTTPException(status_code=400, detail="Qo'llab-quvvatlanmaydigan fayl formati.")

        # Send text straight to the Real Groq API client
        groq_result = analyze_with_real_groq(extracted_text)
        
        if not groq_result["success"]:
            raise HTTPException(status_code=406, detail=groq_result["message"])

        return {
            "status": "success",
            "filename": file.filename,
            "message": "Fayl tasdiqlandi. Groq AI testlarni yaratdi.",
            "data": groq_result
        }
    except HTTPException as he:
        raise he
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/ai/generate-magic")
async def generate_magic_content(payload: dict):
    # Bu route LLM (Gemini/OpenAI) bilan bog'lanib, quiz va flashcard tuzadi
    # Hozircha Mock xizmat ko'rsatadi.
    return {
        "status": "success",
        "flashcards": [
            {"q": "Aritmiya nima?", "a": "Yurak urish ritmining buzilishi."}
        ],
        "quizzes": [
            {
                "question": "Qaysi biri EKGda miokard infarktini anglatadi?",
                "options": ["ST elevatsiyasi", "P tishchasi yo'qolishi", "QRS ingichkaligi"],
                "answer": "ST elevatsiyasi",
                "explanation": "Miokard shikastlanishi darhol ST segment ko'tarilishiga olib keladi."
            }
        ],
        "suggested_images": ["https://example.com/ekg_infarct.jpg"]
    }

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8001)
